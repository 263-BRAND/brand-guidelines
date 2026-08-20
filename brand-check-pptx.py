#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
brand-check-pptx.py — 263 品牌合规检查（PPTX 生成后闸门，可选工具）

Usage:
    python brand-check-pptx.py <file.pptx> [--scheme group-red] [--external]

检查产出 .pptx 是否遵守品牌规范（SKILL.md「PPTX 生成后自查（硬门禁）」）：
  1. 颜色全部来自色板白名单（colorSchemes[scheme] 9 色 + semantic + chartPalette），无自造 hex
  2. 红底/红卡片（品牌红/色阶深档 s4-s7 填充）上的文字必须为白 #FFFFFF，未用 accent 浅粉/浅档/色板外浅色
  3. 结尾页是最后一页（居中 Logo + slogan 图片，无文字；原「感谢/谢谢」页未替换则拦截）
  4. 图片全部内嵌二进制（无外部链接图片）
  5. --external 时全文件不含微软雅黑（对外展示走开源栈）

环境无 python-pptx 时禁止降级——按 SKILL.md 清单人工逐条自查（本脚本是可选机器兜底）。

exit 0 = 通过；exit 1 = 违规（fail-loud）。
"""

import argparse
import sys
import os

try:
    from pptx.oxml.ns import qn
except ImportError:
    qn = None

if hasattr(sys.stdout, 'reconfigure'):
    try:
        sys.stdout.reconfigure(encoding='utf-8', errors='replace')
        sys.stderr.reconfigure(encoding='utf-8', errors='replace')
    except Exception:
        pass

def load_tokens():
    """从脚本所在目录读取 brand-tokens.json（zip 自包含：脚本与 token 同目录）。"""
    script_dir = os.path.dirname(os.path.abspath(__file__))
    for base in (script_dir, os.getcwd()):
        p = os.path.join(base, 'brand-tokens.json')
        if os.path.exists(p):
            import json
            with open(p, 'r', encoding='utf-8') as f:
                return json.load(f)
    sys.exit('brand-check-pptx.py: 未找到 brand-tokens.json（需与脚本同目录或当前目录）。')

# 与 generate.js buildColorWhitelist 同构：colorSchemes[scheme] 全部色值 + semantic 子容器 + chartPalette(series+muted)
META_KEYS = {'note'}
def build_color_whitelist(tokens, scheme):
    ws = {}
    def add(v):
        if isinstance(v, str) and len(v) == 7 and v[0] == '#' and all(c in '0123456789abcdefABCDEF' for c in v[1:]):
            ws[v.lower()] = 1
    cs = tokens.get('colorSchemes', {}).get(scheme)
    if not cs:
        sys.exit('brand-check-pptx.py: 无效 scheme "%s"（brand-tokens.json 无该配色）。' % scheme)
    for k, v in cs.items():
        if k in META_KEYS:
            continue
        if isinstance(v, dict):  # semantic 子容器（positive/negative/neutral），跳过内部 note
            for sk, sv in v.items():
                if sk in META_KEYS:
                    continue
                add(sv)
        else:
            add(v)
    cp = tokens.get('chartPalette', {}).get(scheme)
    if cp:
        for seg in ('series', 'muted'):  # 对齐 generate.js：不加 distinctReds（其色值已在 series 内）
            for item in cp.get(seg, []):
                if isinstance(item, dict):
                    add(item.get('hex'))
    return ws

def rgb_str(rgb):
    """RGBColor → '#rrggbb' 小写。"""
    return '#%02x%02x%02x' % (rgb[0], rgb[1], rgb[2])

# 红底场景判定：品牌红/色阶深档（s4 #EC2831 / s5 #D0121B / s6 #A10D14 / s7 #72090E）
RED_BOTTOM_HEXES = {'#ec2831', '#d0121b', '#a10d14', '#72090e'}

errors = []   # (slide_idx, message) — 违规
warnings = [] # (slide_idx, message) — 无法精确判定，人工确认

def report(slide_idx, msg):
    errors.append((slide_idx, msg))

def warn(slide_idx, msg):
    warnings.append((slide_idx, msg))

def check_color(slide_idx, what, hexval, whitelist):
    """白名单检查（红底系判定单独走 is_red_bottom_hex）。"""
    hexval = hexval.lower()
    if hexval not in whitelist:
        report(slide_idx, '%s 颜色 %s 不在品牌色板白名单内（SKILL.md「品牌规则 → 色彩」）。' % (what, hexval))

def check_text_runs(slide_idx, label, text_frame, whitelist, on_red=False):
    """遍历段落 run：文字色白名单 + 红底上文字必须白。on_red 由几何相交判定注入。"""
    for para in text_frame.paragraphs:
        for run in para.runs:
            try:
                c = run.font.color
                if c and c.type is not None and getattr(c, 'rgb', None) is not None:
                    hexv = rgb_str(c.rgb)
                    check_color(slide_idx, label + '文字', hexv, whitelist)
                    if on_red and hexv.lower() != '#ffffff':
                        report(slide_idx, '%s 位于红底（品牌红/色阶深档）上，文字颜色 %s 非白 #FFFFFF——红底次级元素必须白色，禁止 accent 浅粉/浅档（红底+浅粉=死色）。' % (label, hexv))
            except Exception:
                pass  # 主题色/无法读取 → 人工确认层处理
            # 中文字体可能设在 a:ea 面（run.font.name 只读 latin 面）——latin + ea 都查
            fnames = []
            try:
                if run.font.name:
                    fnames.append(run.font.name)
                if qn is not None and run._r is not None:
                    rPr = run._r.get_or_add_rPr()
                    ea = rPr.find(qn('a:ea'))
                    if ea is not None and ea.get('typeface'):
                        fnames.append(ea.get('typeface'))
            except Exception:
                pass
            if external and any('微软雅黑' in f or 'microsoft yahei' in f or 'microsoftyahei' in f for f in fnames):
                report(slide_idx, '%s 使用微软雅黑（闭源），对外展示必须用开源栈 Noto Sans SC/Source Han Sans SC。' % label)

def shape_bounds(shape):
    """shape → (left, top, right, bottom) EMU；无几何数据返回 None。"""
    try:
        if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
            return None
        return (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)
    except Exception:
        return None

def rects_overlap(a, b):
    """两矩形相交（含边界相接）。a/b 为 (l,t,r,b)。"""
    return not (a[2] <= b[0] or b[2] <= a[0] or a[3] <= b[1] or b[3] <= a[1])

def is_red_bottom_hex(hexv):
    return hexv.lower() in RED_BOTTOM_HEXES

def check_slide_shapes(slide_idx, shapes, whitelist):
    """两遍扫描：先收集红底形状几何框，再检查文字框（几何相交 → 红底上文字必须白）。"""
    # 遍 1：收集红底填充形状的边界框 + 校验形状填充色白名单
    red_boxes = []
    for shape in shapes:
        try:
            if shape.fill.type == 1:  # MSO_FILL_TYPE.SOLID
                fc = shape.fill.fore_color
                if fc.type is not None and getattr(fc, 'rgb', None) is not None:
                    hexv = rgb_str(fc.rgb)
                    check_color(slide_idx, '形状填充', hexv, whitelist)
                    if is_red_bottom_hex(hexv):
                        b = shape_bounds(shape)
                        if b:
                            red_boxes.append(b)
        except Exception:
            pass

    # 遍 2：文字框/表格 → 与任一红底框相交则该处文字必须白
    for shape in shapes:
        tb = shape_bounds(shape)
        on_red = bool(red_boxes and tb and any(rects_overlap(tb, rb) for rb in red_boxes))
        if shape.has_text_frame:
            check_text_runs(slide_idx, '文字框', shape.text_frame, whitelist, on_red)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    check_text_runs(slide_idx, '表格', cell.text_frame, whitelist, on_red)

    # 遍 3：图片内嵌 + 图表系列色
    for shape in shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            try:
                shape.image.blob
            except Exception as e:
                report(slide_idx, '图片使用外部路径链接（未内嵌二进制）。python-pptx 无法读取：%s。品牌图片必须内嵌。' % e)
        if shape.has_chart:
            try:
                for series in shape.chart.series:
                    for point in series.points:
                        try:
                            pf = point.format.fill
                            if pf.type == 1:
                                pc = pf.fore_color
                                if pc.type is not None and getattr(pc, 'rgb', None) is not None:
                                    check_color(slide_idx, '图表系列', rgb_str(pc.rgb), whitelist)
                        except Exception:
                            pass
            except Exception:
                pass

def collect_slide_text(shapes):
    """收集页内全部文字（文本框 + 表格单元格），返回去重非空列表。"""
    texts = []
    for s in shapes:
        if s.has_text_frame and s.text_frame.text.strip():
            texts.append(s.text_frame.text.strip())
        if s.has_table:
            for row in s.table.rows:
                for cell in row.cells:
                    t = cell.text_frame.text.strip()
                    if t:
                        texts.append(t)
    return texts

def check_end_slide(slide_idx, shapes, is_last):
    """结尾页规则。is_last=True：最后一页必须 VI 标准（居中 Logo+slogan 图、无文字、无「感谢/谢谢」）；
    is_last=False：中间页出现「感谢/谢谢」结尾字样 → 疑似原文件结尾/感谢页残留（对齐 generate.js「中间夹 end 拦截」）。
    注意：中间页可能正常含「感谢 XX 团队」等致谢句，此处降级为 warning（人工确认），只有最后一页才是硬拦截。"""
    pics = [s for s in shapes if s.shape_type == 13]
    texts = collect_slide_text(shapes)
    all_text = ' '.join(texts).lower()
    thanks = '感谢' in all_text or '谢谢' in all_text or 'thanks' in all_text
    if is_last:
        if not pics:
            report(slide_idx, '结尾页（最后一页）无 Logo/slogan 图片——VI 标准结尾页必须居中 Logo + slogan PNG。')
        if thanks:
            report(slide_idx, '结尾页（最后一页）含「感谢/谢谢」文字——原文件结尾/感谢页未替换为 VI 标准结尾页。')
        elif texts:
            warn(slide_idx, '结尾页（最后一页）含文字：「%s」。VI 标准结尾页应无文字（居中 Logo + slogan）；若有合理说明文字请人工确认。' % ' / '.join(texts[:3]))
    elif thanks:
        warn(slide_idx, '中间页含「感谢/谢谢」结尾字样——疑似原文件结尾/感谢页残留（结尾页必须是整个 PPT 唯一最后一页）。若为内容中的致谢句请人工确认。')

def main():
    global external
    ap = argparse.ArgumentParser(description='263 品牌合规检查（PPTX 生成后闸门）')
    ap.add_argument('file', help='要检查的 .pptx 文件')
    ap.add_argument('--scheme', default='group-red', help='配色方案（brand-tokens.json colorSchemes 键），默认 group-red')
    ap.add_argument('--external', action='store_true', help='对外展示：额外检查全文件不含微软雅黑')
    args = ap.parse_args()
    external = args.external

    if not os.path.exists(args.file):
        sys.exit('brand-check-pptx.py: 文件不存在：%s' % args.file)
    if not args.file.lower().endswith('.pptx'):
        sys.exit('brand-check-pptx.py: 仅支持 .pptx（文件：%s）' % args.file)

    tokens = load_tokens()
    whitelist = build_color_whitelist(tokens, args.scheme)
    if not whitelist:
        sys.exit('brand-check-pptx.py: 白名单为空——brand-tokens.json %s 无色值可查。' % args.scheme)

    try:
        from pptx import Presentation
    except ImportError:
        sys.exit('brand-check-pptx.py: 缺少 python-pptx（本机未安装）。环境无 python-pptx 时按 SKILL.md「PPTX 生成后自查」清单人工逐条自查，不降级约束。')

    prs = Presentation(args.file)
    slides = list(prs.slides)
    n = len(slides)
    if n == 0:
        sys.exit('brand-check-pptx.py: 文件无任何幻灯片。')

    for idx, slide in enumerate(slides):
        shapes = list(slide.shapes)
        check_slide_shapes(idx, shapes, whitelist)
        check_end_slide(idx, shapes, is_last=(idx == n - 1))

    # 汇总
    if errors:
        print('品牌合规检查未通过（%s）：' % os.path.basename(args.file))
        for idx, msg in errors:
            print('  slide %d: %s' % (idx, msg))
        print('按 SKILL.md「PPTX 生成后自查」修复后重跑；不要带病交付。')
        sys.exit(1)
    if warnings:
        print('品牌合规检查通过（有 %d 条需人工确认）：' % len(warnings), os.path.basename(args.file))
        for idx, msg in warnings:
            print('  slide %d: %s' % (idx, msg))
    else:
        print('品牌合规检查通过：%s' % os.path.basename(args.file))
    sys.exit(0)

if __name__ == '__main__':
    main()
